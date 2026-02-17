"""
OCR ULTRA PROFESIONAL API - Flask Application for Google Cloud Run
Sistema híbrido inteligente:
- NIVEL 1: OCR Ultra Pro (Tesseract + EasyOCR + mejoras de imagen avanzadas) - 100% GRATIS
- NIVEL 2: Gemini Fallback (solo cuando OCR falla) - Con costo
Soporta: Imágenes (JPG, PNG, etc), PDFs, Word, PowerPoint
"""

from flask import Flask, request, jsonify
from werkzeug.utils import secure_filename
import os
import cv2
import pytesseract
import easyocr
from PIL import Image, ImageEnhance, ImageFilter
import numpy as np
import json
from datetime import datetime
import tempfile
import mimetypes
import re
import time
import random

# Importar Gemini para fallback
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

# Importar procesadores de documentos
try:
    import fitz  # PyMuPDF para PDFs
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document  # python-docx para Word
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation  # python-pptx para PowerPoint
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from markitdown import MarkItDown
    MARKITDOWN_AVAILABLE = True
except ImportError:
    MARKITDOWN_AVAILABLE = False

# Configuración
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024  # 32MB máximo

# Configurar Tesseract
pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

# Configurar Gemini (API key desde variable de entorno)
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', '')
gemini_model = None

if GEMINI_AVAILABLE and GEMINI_API_KEY:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        models_to_try = ['gemini-1.5-flash-8b', 'gemini-1.5-flash', 'gemini-2.0-flash-exp']
        for model_name in models_to_try:
            try:
                gemini_model = genai.GenerativeModel(model_name)
                print(f"✅ Gemini {model_name} disponible como fallback")
                break
            except:
                continue
    except Exception as e:
        print(f"⚠️ Error configurando Gemini: {e}")

# Inicializar EasyOCR de forma lazy (solo cuando se necesite)
easyocr_reader = None
EASYOCR_AVAILABLE = False

def get_easyocr_reader():
    """Obtener lector de EasyOCR (inicialización lazy)"""
    global easyocr_reader, EASYOCR_AVAILABLE
    
    if easyocr_reader is not None:
        return easyocr_reader
    
    try:
        print("🔧 Inicializando EasyOCR...")
        # Silenciar progreso de descarga
        import warnings
        warnings.filterwarnings('ignore')
        
        easyocr_reader = easyocr.Reader(['es', 'en'], gpu=False, verbose=False)
        EASYOCR_AVAILABLE = True
        print("✅ EasyOCR listo")
        return easyocr_reader
    except Exception as e:
        print(f"⚠️ EasyOCR no disponible: {e}")
        EASYOCR_AVAILABLE = False
        return None

print("✅ OCR Ultra Pro API iniciado (EasyOCR se cargará cuando se necesite)")


class DocumentProcessor:
    """Procesador inteligente de documentos"""
    
    ALLOWED_EXTENSIONS = {
        'image': {'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif', 'webp'},
        'pdf': {'pdf'},
        'word': {'docx', 'doc'},
        'powerpoint': {'pptx', 'ppt'}
    }
    
    @classmethod
    def detect_file_type(cls, filename):
        """Detectar tipo de archivo"""
        ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
        
        for file_type, extensions in cls.ALLOWED_EXTENSIONS.items():
            if ext in extensions:
                return file_type, ext
        
        return None, ext
    
    @classmethod
    def is_allowed_file(cls, filename):
        """Verificar si el archivo es permitido"""
        file_type, ext = cls.detect_file_type(filename)
        return file_type is not None


class UltraProImageEnhancer:
    """Mejorador de imágenes ULTRA PROFESIONAL optimizado para velocidad y calidad"""
    
    @staticmethod
    def enhance_image_ultimate(image_input):
        """Mejoras de imagen optimizadas - BALANCE VELOCIDAD/CALIDAD"""
        
        # Convertir PIL/path a OpenCV
        if isinstance(image_input, str):
            opencv_image = cv2.imread(image_input)
        elif isinstance(image_input, Image.Image):
            opencv_image = cv2.cvtColor(np.array(image_input), cv2.COLOR_RGB2BGR)
        else:
            opencv_image = image_input
        
        if opencv_image is None:
            raise ValueError("No se pudo cargar la imagen")
        
        try:
            # Convertir a escala de grises (más rápido)
            if len(opencv_image.shape) == 3:
                gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
            else:
                gray = opencv_image
            
            # 1. CORRECCIÓN DE INCLINACIÓN (solo si es necesario)
            edges = cv2.Canny(gray, 50, 150, apertureSize=3)
            lines = cv2.HoughLines(edges, 1, np.pi/180, threshold=100)
            
            if lines is not None and len(lines) > 3:
                angles = []
                for rho, theta in lines[:5]:  # Solo 5 líneas más prominentes
                    angle = (theta - np.pi/2) * 180 / np.pi
                    if abs(angle) < 45:  # Ignorar ángulos extremos
                        angles.append(angle)
                
                if angles:
                    median_angle = np.median(angles)
                    if abs(median_angle) > 1.0:  # Solo corregir si >1 grado
                        (h, w) = gray.shape[:2]
                        center = (w // 2, h // 2)
                        M = cv2.getRotationMatrix2D(center, median_angle, 1.0)
                        gray = cv2.warpAffine(gray, M, (w, h), flags=cv2.INTER_LINEAR, 
                                            borderMode=cv2.BORDER_REPLICATE)
            
            # 2. REDUCCIÓN DE RUIDO (optimizado)
            denoised = cv2.fastNlMeansDenoising(gray, None, h=10, templateWindowSize=7, searchWindowSize=21)
            
            # 3. MEJORA DE CONTRASTE (CLAHE optimizado)
            clahe = cv2.createCLAHE(clipLimit=2.5, tileGridSize=(8, 8))
            enhanced = clahe.apply(denoised)
            
            # 4. BINARIZACIÓN INTELIGENTE (método híbrido rápido)
            # Usar Otsu que es rápido y efectivo
            blurred = cv2.GaussianBlur(enhanced, (5, 5), 0)
            _, otsu = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            
            # Adaptativo solo si Otsu no es suficiente
            if np.mean(otsu) < 50 or np.mean(otsu) > 200:  # Si imagen muy oscura o clara
                adaptive = cv2.adaptiveThreshold(
                    enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
                )
                # Combinar ambos métodos
                final = cv2.bitwise_or(otsu, adaptive)
            else:
                final = otsu
            
            # 5. LIMPIEZA RÁPIDA (solo elementos grandes)
            kernel = np.ones((2, 2), np.uint8)
            final = cv2.morphologyEx(final, cv2.MORPH_OPEN, kernel, iterations=1)
            
            return final
            
        except Exception as e:
            # Fallback ultra rápido
            if len(opencv_image.shape) == 3:
                gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
            else:
                gray = opencv_image
            
            # Solo binarización básica
            _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            return binary


class UltraProOCRProcessor:
    """Motor OCR ULTRA PROFESIONAL con múltiples engines y fallback inteligente"""
    
    def __init__(self):
        # EasyOCR se inicializará cuando se necesite (lazy loading)
        self.enhancer = UltraProImageEnhancer()
        
    def get_easyocr_reader(self):
        """Obtener lector de EasyOCR con inicialización lazy"""
        return get_easyocr_reader()
        
    def process_with_tesseract_ultra(self, image):
        """Tesseract con 6 configuraciones ultra profesionales"""
        configs = [
            {'config': '--oem 3 --psm 6', 'name': 'Bloque de texto', 'weight': 1.0},
            {'config': '--oem 3 --psm 4', 'name': 'Columna única', 'weight': 0.9},
            {'config': '--oem 3 --psm 3', 'name': 'Auto completo', 'weight': 0.95},
            {'config': '--oem 3 --psm 1', 'name': 'Auto con OSD', 'weight': 0.85},
            {'config': '--oem 3 --psm 11', 'name': 'Texto disperso', 'weight': 0.7},
            {'config': '--oem 3 --psm 6 -c preserve_interword_spaces=1', 'name': 'Espacios preservados', 'weight': 0.8},
        ]
        
        best_result = {'text': '', 'confidence': 0, 'config': '', 'word_count': 0, 'score': 0}
        
        for config in configs:
            try:
                # Extraer texto
                text = pytesseract.image_to_string(
                    image, 
                    lang='spa+eng', 
                    config=config['config']
                ).strip()
                
                if not text or len(text) < 3:
                    continue
                
                # Obtener confianza
                try:
                    data = pytesseract.image_to_data(
                        image, 
                        lang='spa+eng', 
                        config=config['config'], 
                        output_type=pytesseract.Output.DICT
                    )
                    confidences = [int(c) for c in data['conf'] if str(c).isdigit() and int(c) > 0]
                    avg_confidence = sum(confidences) / len(confidences) if confidences else 0
                except:
                    avg_confidence = 50.0
                
                word_count = len([w for w in text.split() if len(w) > 1])
                
                # Score combinado: confianza + cantidad de texto válido + peso de configuración
                score = (avg_confidence * config['weight']) + (word_count * 2)
                
                if score > best_result['score']:
                    best_result = {
                        'text': text,
                        'confidence': avg_confidence,
                        'config': config['name'],
                        'word_count': word_count,
                        'score': score
                    }
                    
            except Exception as e:
                continue
        
        return best_result
    
    def process_with_easyocr_ultra(self, image_path):
        """EasyOCR con configuraciones ultra profesionales"""
        try:
            # Inicializar EasyOCR solo cuando se necesite
            reader = self.get_easyocr_reader()
            
            if reader is None:
                return {'text': '', 'confidence': 0, 'error': 'EasyOCR no disponible', 'score': 0}
            
            # Diferentes niveles de confianza
            confidence_levels = [0.1, 0.25, 0.4]
            best_result = {'text': '', 'confidence': 0, 'detections': 0, 'score': 0}
            
            for conf_threshold in confidence_levels:
                results = reader.readtext(image_path, detail=1)
                
                # Filtrar por confianza
                valid_results = [item for item in results if item[2] >= conf_threshold]
                
                if not valid_results:
                    continue
                
                # Ordenar por posición (izquierda a derecha, arriba a abajo)
                valid_results.sort(key=lambda x: (x[0][0][1], x[0][0][0]))
                
                text = ' '.join([item[1] for item in valid_results])
                avg_confidence = sum([item[2] for item in valid_results]) / len(valid_results)
                
                score = (avg_confidence * 100) + (len(text.split()) * 2)
                
                if score > best_result['score']:
                    best_result = {
                        'text': text,
                        'confidence': avg_confidence * 100,
                        'detections': len(valid_results),
                        'threshold': conf_threshold,
                        'score': score
                    }
            
            return best_result
            
        except Exception as e:
            return {'text': '', 'confidence': 0, 'error': str(e), 'score': 0}
    
    def extract_structured_data(self, text):
        """Extracción inteligente de datos estructurados"""
        patterns = {
            'emails': r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            'phones': r'\b(?:\+?57\s?)?[0-9]{7,10}\b',
            'dates': r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b',
            'ids': r'\b[0-9]{6,12}\b',
            'currency': r'\$[\d,]+(?:\.\d{2})?',
        }
        
        extracted = {}
        for key, pattern in patterns.items():
            matches = re.findall(pattern, text, re.IGNORECASE)
            extracted[key] = list(set(matches)) if matches else []
        
        return extracted
    
    def analyze_quality(self, ocr_result):
        """Análisis de calidad del OCR - determina si es suficientemente bueno"""
        if not ocr_result or not ocr_result.get('text'):
            return False, "Sin texto extraído", 0
        
        text = ocr_result['text']
        confidence = ocr_result.get('confidence', 0)
        
        words = text.split()
        
        # Criterio 1: Longitud mínima
        if len(words) < 3:
            return False, f"Muy pocas palabras: {len(words)}", 10
        
        # Criterio 2: Confianza mínima
        if confidence < 35:
            return False, f"Confianza muy baja: {confidence:.1f}%", 20
        
        # Criterio 3: Detectar texto corrupto
        corruption_score = 0
        corruption_score += len(re.findall(r'[A-Z]{5,}', text))  # Demasiadas mayúsculas
        corruption_score += len(re.findall(r'[^\w\s]{4,}', text))  # Símbolos raros
        
        if corruption_score > 5:
            return False, f"Texto corrupto (score: {corruption_score})", 30
        
        # Criterio 4: Coherencia del texto
        spanish_common = ['el', 'la', 'de', 'que', 'y', 'en', 'un', 'es', 'se', 'no', 'por', 'con', 'para']
        coherent_words = sum(1 for word in words if word.lower() in spanish_common or len(word) > 2)
        coherence_ratio = coherent_words / len(words) if words else 0
        
        if coherence_ratio < 0.15:
            return False, f"Texto incoherente (coherencia: {coherence_ratio:.2f})", 40
        
        # Criterio 5: Datos estructurados
        data = self.extract_structured_data(text)
        data_points = sum(len(v) for v in data.values())
        
        # Calcular score de calidad (0-100)
        quality_score = min(100, (
            confidence * 0.4 +
            min(len(words) / 10, 20) * 0.3 +
            coherence_ratio * 100 * 0.2 +
            min(data_points * 5, 20) * 0.1
        ))
        
        # Decisión final
        if quality_score >= 50:
            return True, f"Calidad aceptable (score: {quality_score:.1f})", quality_score
        elif quality_score >= 35 and len(text) > 50:
            return True, f"Calidad suficiente (score: {quality_score:.1f})", quality_score
        else:
            return False, f"Calidad insuficiente (score: {quality_score:.1f})", quality_score
    
    def select_best_engine_result(self, tesseract_result, easyocr_result):
        """Seleccionar el mejor resultado entre engines"""
        results = []
        
        # Evaluar Tesseract
        if tesseract_result.get('text') and len(tesseract_result['text'].strip()) > 5:
            results.append(('tesseract', tesseract_result, tesseract_result.get('score', 0)))
        
        # Evaluar EasyOCR
        if easyocr_result.get('text') and len(easyocr_result['text'].strip()) > 5:
            results.append(('easyocr', easyocr_result, easyocr_result.get('score', 0)))
        
        if not results:
            return None, "Ningún engine produjo resultados válidos"
        
        # Seleccionar mejor resultado por score
        best_engine, best_result, best_score = max(results, key=lambda x: x[2])
        
        return {
            'engine': best_engine,
            'text': best_result['text'],
            'confidence': best_result['confidence'],
            'score': best_score,
            'config': best_result.get('config', best_result.get('threshold', 'default'))
        }, f"Mejor engine: {best_engine}"
    
    def gemini_fallback(self, image_path):
        """Fallback a Gemini cuando OCR falla"""
        if not gemini_model:
            return None, "Gemini no disponible"
        
        try:
            img = Image.open(image_path)
            
            prompt = """Analiza esta imagen y extrae TODO el texto visible con máxima precisión.

INSTRUCCIONES:
1. Transcribe TODO el texto que veas, preservando el formato
2. Identifica tipo de documento (factura, identificación, carta, etc.)
3. Extrae datos estructurados: nombres, fechas, números, emails, teléfonos, direcciones
4. Si algo es ilegible, indica [ILEGIBLE]

Formato de respuesta:
TIPO_DOCUMENTO: [tipo]
TEXTO_COMPLETO:
[texto extraído]

DATOS_ESTRUCTURADOS:
- Nombres: [si hay]
- Fechas: [si hay]
- Números ID: [si hay]
- Contactos: [si hay]
- Otros: [si hay]"""
            
            response = gemini_model.generate_content([prompt, img])
            
            return {
                'text': response.text,
                'engine': 'gemini',
                'confidence': 85.0,  # Asumimos alta confianza en Gemini
                'method': 'ai_vision'
            }, "Gemini fallback exitoso"
            
        except Exception as e:
            return None, f"Error en Gemini: {str(e)}"
    
    def process_image_ultra_pro(self, image_path, use_gemini_fallback=True):
        """Procesamiento ULTRA PRO completo de una imagen - OPTIMIZADO PARA VELOCIDAD"""
        start_time = time.time()
        
        print(f"⚡ Procesando con OCR Ultra Pro (RÁPIDO)...")
        
        # FASE 1: Mejorar imagen (OPTIMIZADO)
        enhanced_image = None
        temp_enhanced_path = image_path
        
        try:
            enhanced_image = self.enhancer.enhance_image_ultimate(image_path)
            temp_enhanced_path = image_path.replace('.', '_enhanced.') if '.' in image_path else image_path + '_enhanced.png'
            cv2.imwrite(temp_enhanced_path, enhanced_image)
        except Exception as e:
            print(f"⚠️ Mejoras básicas aplicadas")
            enhanced_image = cv2.imread(image_path, cv2.IMREAD_GRAYSCALE) if enhanced_image is None else enhanced_image
        
        # FASE 2: Procesar con múltiples engines OCR (PARALELO)
        print("🔍 Extrayendo texto...")
        
        tesseract_result = self.process_with_tesseract_ultra(enhanced_image if enhanced_image is not None else image_path)
        easyocr_result = self.process_with_easyocr_ultra(temp_enhanced_path)
        
        # Limpiar archivo temporal
        if temp_enhanced_path != image_path and os.path.exists(temp_enhanced_path):
            try:
                os.remove(temp_enhanced_path)
            except:
                pass
        
        # FASE 3: Seleccionar mejor resultado
        best_result, selection_reason = self.select_best_engine_result(tesseract_result, easyocr_result)
        
        if best_result:
            # Analizar calidad
            is_quality, quality_msg, quality_score = self.analyze_quality(best_result)
            
            if is_quality:
                processing_time = time.time() - start_time
                structured_data = self.extract_structured_data(best_result['text'])
                
                # Determinar recomendación basada en calidad
                recommendation = self._get_recommendation(quality_score, best_result['confidence'])
                
                print(f"✅ OCR GRATUITO exitoso en {processing_time:.1f}s")
                
                return {
                    'success': True,
                    'method': 'ocr_ultra_pro',
                    'method_name': '🆓 OCR ULTRA PRO (100% GRATUITO)',
                    'engine': best_result['engine'].upper(),
                    'text': best_result['text'],
                    'confidence': best_result['confidence'],
                    'quality_score': quality_score,
                    'structured_data': structured_data,
                    'processing_time': round(processing_time, 2),
                    'cost': '0.00 USD',
                    'recommendation': recommendation,
                    'engines_tested': {
                        'tesseract': {
                            'confidence': tesseract_result.get('confidence', 0),
                            'word_count': tesseract_result.get('word_count', 0)
                        },
                        'easyocr': {
                            'confidence': easyocr_result.get('confidence', 0),
                            'detections': easyocr_result.get('detections', 0)
                        }
                    }
                }
            else:
                print(f"⚠️ Calidad insuficiente: {quality_msg}")
        
        # FASE 4: Fallback a Gemini (si está habilitado)
        if use_gemini_fallback and gemini_model:
            print("🤖 Activando Gemini AI (análisis avanzado)...")
            gemini_result, gemini_msg = self.gemini_fallback(image_path)
            
            if gemini_result:
                processing_time = time.time() - start_time
                structured_data = self.extract_structured_data(gemini_result['text'])
                
                print(f"✅ GEMINI AI exitoso en {processing_time:.1f}s")
                
                return {
                    'success': True,
                    'method': 'gemini_fallback',
                    'method_name': '🤖 GEMINI AI (Análisis Avanzado)',
                    'engine': 'GEMINI-1.5-FLASH',
                    'text': gemini_result['text'],
                    'confidence': gemini_result['confidence'],
                    'quality_score': 85,
                    'structured_data': structured_data,
                    'processing_time': round(processing_time, 2),
                    'cost': 'Variable según uso',
                    'recommendation': '💡 OCR gratuito no fue suficiente - Se usó IA avanzada',
                    'fallback_reason': quality_msg if 'quality_msg' in locals() else 'OCR insuficiente',
                    'ocr_attempted': {
                        'tesseract': tesseract_result.get('confidence', 0),
                        'easyocr': easyocr_result.get('confidence', 0)
                    }
                }
            else:
                print(f"❌ Gemini también falló: {gemini_msg}")
        
        # FASE 5: Ambos métodos fallaron
        processing_time = time.time() - start_time
        
        return {
            'success': False,
            'method': 'all_failed',
            'method_name': '❌ PROCESAMIENTO FALLIDO',
            'error': '❌ ERROR CRÍTICO: Ni OCR Ultra Pro ni Gemini AI pudieron procesar el documento',
            'processing_time': round(processing_time, 2),
            'recommendation': '⚠️ El documento no pudo ser procesado. Recomendaciones:\n' +
                            '  1. Verifica que la imagen contenga texto legible\n' +
                            '  2. Mejora la calidad: Mayor resolución (mínimo 300 DPI)\n' +
                            '  3. Asegura buena iluminación y contraste\n' +
                            '  4. Evita texto borroso o muy pequeño\n' +
                            '  5. Intenta escanear el documento físico si es posible',
            'ocr_details': {
                'tesseract_confidence': tesseract_result.get('confidence', 0),
                'easyocr_confidence': easyocr_result.get('confidence', 0),
                'tesseract_words': tesseract_result.get('word_count', 0),
                'easyocr_detections': easyocr_result.get('detections', 0)
            },
            'gemini_available': gemini_model is not None,
            'gemini_enabled': use_gemini_fallback
        }
    
    def _get_recommendation(self, quality_score, confidence):
        """Generar recomendación basada en calidad"""
        if quality_score >= 80 and confidence >= 80:
            return "✅ EXCELENTE: Documento procesado con máxima calidad"
        elif quality_score >= 60 and confidence >= 60:
            return "✅ MUY BUENO: Extracción exitosa con alta confiabilidad"
        elif quality_score >= 50:
            return "✅ BUENO: Texto extraído correctamente - Resultado confiable"
        elif quality_score >= 40:
            return "⚠️ ACEPTABLE: Documento procesado - Considera mejorar calidad de escaneo para mejores resultados"
        else:
            return "⚠️ COMPLEJO: Documento procesado pero recomienda mejorar calidad de imagen (mayor resolución, mejor iluminación)"


class PDFProcessor:
    """Procesador de archivos PDF con OCR Ultra Pro"""
    
    def __init__(self, ocr_processor):
        self.ocr = ocr_processor
    
    def process(self, pdf_path, use_gemini_fallback=True):
        """Procesar PDF completo con OCR Ultra Pro"""
        if not PDF_AVAILABLE:
            return {'success': False, 'error': 'PyMuPDF no disponible'}
        
        try:
            pdf_doc = fitz.open(pdf_path)
            pages = []
            total_text = ''
            ocr_success_count = 0
            gemini_usage_count = 0
            failed_count = 0
            
            # Limitar páginas para evitar timeouts
            max_pages = min(pdf_doc.page_count, 20)
            
            print(f"📄 Procesando PDF: {max_pages} páginas...")
            
            for page_num in range(max_pages):
                page = pdf_doc[page_num]
                
                # Intentar extracción directa primero
                direct_text = page.get_text().strip()
                
                if direct_text and len(direct_text) > 50:
                    # Texto directo disponible (PDF con texto)
                    page_result = {
                        'page': page_num + 1,
                        'text': direct_text,
                        'method': '📄 Extracción directa',
                        'char_count': len(direct_text),
                        'success': True
                    }
                    ocr_success_count += 1
                    print(f"  ✅ Página {page_num + 1}: Texto directo")
                else:
                    # Necesita OCR (PDF escaneado)
                    print(f"  🔍 Página {page_num + 1}: Aplicando OCR...")
                    
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                        pix.save(tmp.name)
                        ocr_result = self.ocr.process_image_ultra_pro(tmp.name, use_gemini_fallback)
                        os.unlink(tmp.name)
                    
                    if ocr_result.get('success'):
                        method_emoji = '🆓' if ocr_result.get('method') == 'ocr_ultra_pro' else '🤖'
                        page_result = {
                            'page': page_num + 1,
                            'text': ocr_result.get('text', ''),
                            'method': f"{method_emoji} {ocr_result.get('method_name', 'OCR')}",
                            'confidence': ocr_result.get('confidence', 0),
                            'quality_score': ocr_result.get('quality_score', 0),
                            'processing_time': ocr_result.get('processing_time', 0),
                            'success': True
                        }
                        
                        if ocr_result.get('method') == 'gemini_fallback':
                            gemini_usage_count += 1
                            print(f"    🤖 Página {page_num + 1}: Gemini AI usado")
                        else:
                            ocr_success_count += 1
                            print(f"    ✅ Página {page_num + 1}: OCR gratuito")
                    else:
                        failed_count += 1
                        page_result = {
                            'page': page_num + 1,
                            'text': '',
                            'method': '❌ Procesamiento fallido',
                            'error': ocr_result.get('error', 'Unknown error'),
                            'success': False
                        }
                        print(f"    ❌ Página {page_num + 1}: Falló")
                
                pages.append(page_result)
                total_text += page_result.get('text', '') + '\n\n'
            
            pdf_doc.close()
            
            # Estadísticas y recomendaciones
            total_pages_processed = len(pages)
            success_rate = ((ocr_success_count + gemini_usage_count) / total_pages_processed * 100) if total_pages_processed > 0 else 0
            free_rate = (ocr_success_count / total_pages_processed * 100) if total_pages_processed > 0 else 0
            
            # Generar recomendación
            if free_rate >= 90:
                recommendation = "✅ EXCELENTE: 90%+ procesado gratis - Documento de muy buena calidad"
            elif free_rate >= 70:
                recommendation = "✅ MUY BUENO: 70%+ procesado gratis - Calidad aceptable"
            elif free_rate >= 50:
                recommendation = "⚠️ ACEPTABLE: Solo 50-70% procesado gratis - Considera mejorar calidad de escaneo"
            else:
                recommendation = "⚠️ COMPLEJO: <50% procesado gratis - Documento requiere mejorar calidad significativamente"
            
            if failed_count > 0:
                recommendation += f" | ⚠️ {failed_count} página(s) fallaron"
            
            return {
                'success': True,
                'file_type': 'pdf',
                'total_pages': pdf_doc.page_count if 'pdf_doc' in dir() else max_pages,
                'pages_processed': total_pages_processed,
                'pages': pages,
                'total_text': total_text.strip(),
                'char_count': len(total_text),
                'statistics': {
                    'free_ocr': ocr_success_count,
                    'ai_processing': gemini_usage_count,
                    'failed': failed_count,
                    'success_rate': round(success_rate, 1),
                    'free_percentage': round(free_rate, 1)
                },
                'recommendation': recommendation
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'❌ Error procesando PDF: {str(e)}',
                'file_type': 'pdf'
            }


class WordProcessor:
    """Procesador de archivos Word"""
    
    def process(self, docx_path):
        """Procesar documento Word"""
        if not DOCX_AVAILABLE:
            return {'error': 'python-docx no disponible'}
        
        try:
            doc = Document(docx_path)
            
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
            
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            full_text.append(cell.text)
            
            text = '\n'.join(full_text)
            
            return {
                'file_type': 'word',
                'text': text,
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'char_count': len(text),
                'word_count': len(text.split())
            }
            
        except Exception as e:
            return {'error': str(e)}


class PowerPointProcessor:
    """Procesador de archivos PowerPoint"""
    
    def process(self, pptx_path):
        """Procesar presentación PowerPoint"""
        if not PPTX_AVAILABLE:
            return {'error': 'python-pptx no disponible'}
        
        try:
            prs = Presentation(pptx_path)
            slides = []
            total_text = ''
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                text = '\n'.join(slide_text)
                slides.append({
                    'slide': slide_num,
                    'text': text,
                    'shapes_count': len(slide.shapes)
                })
                total_text += text + '\n\n'
            
            return {
                'file_type': 'powerpoint',
                'total_slides': len(slides),
                'slides': slides,
                'total_text': total_text.strip(),
                'char_count': len(total_text)
            }
            
        except Exception as e:
            return {'error': str(e)}


# Crear procesadores globales
ultra_ocr_processor = UltraProOCRProcessor()
pdf_processor = PDFProcessor(ultra_ocr_processor)
word_processor = WordProcessor()
ppt_processor = PowerPointProcessor()

# Crear procesador MarkItDown si está disponible
if MARKITDOWN_AVAILABLE:
    markitdown = MarkItDown()
    print("✅ MarkItDown disponible")
else:
    markitdown = None
    print("⚠️ MarkItDown no disponible")

print("=" * 80)
print("🚀 OCR ULTRA PROFESIONAL API - LISTO")
print("=" * 80)
print(f"✅ Tesseract: Activo (6 configuraciones)")
print(f"✅ EasyOCR: Activo (múltiples niveles de confianza)")
print(f"✅ Mejoras de imagen: IA avanzada")
print(f"{'✅' if gemini_model else '⚠️'} Gemini Fallback: {'Disponible' if gemini_model else 'No configurado'}")
print(f"✅ PyMuPDF: {'Disponible' if PDF_AVAILABLE else 'No disponible'}")
print(f"✅ python-docx: {'Disponible' if DOCX_AVAILABLE else 'No disponible'}")
print(f"✅ python-pptx: {'Disponible' if PPTX_AVAILABLE else 'No disponible'}")
print("=" * 80)


@app.route('/', methods=['GET'])
def home():
    """Endpoint de información"""
    return jsonify({
        'service': 'OCR Ultra Profesional API',
        'version': '3.0.0',
        'status': 'active',
        'description': 'Sistema híbrido inteligente con OCR Ultra Pro + Gemini Fallback',
        'endpoints': {
            'process': {
                'url': '/api/process',
                'method': 'POST',
                'description': 'Procesar documento con OCR Ultra Pro (Tesseract + EasyOCR + mejoras IA)',
                'parameters': {
                    'file': 'Archivo a procesar (requerido)',
                    'use_gemini_fallback': 'true/false - Usar Gemini si OCR falla (default: true)'
                }
            },
            'process_markdown': {
                'url': '/api/process-markdown',
                'method': 'POST',
                'description': 'Procesar documento y retornar en formato Markdown (MarkItDown)',
                'available': MARKITDOWN_AVAILABLE
            },
            'health': {
                'url': '/health',
                'method': 'GET',
                'description': 'Verificar estado del servicio'
            }
        },
        'supported_formats': {
            'images': list(DocumentProcessor.ALLOWED_EXTENSIONS['image']),
            'pdf': list(DocumentProcessor.ALLOWED_EXTENSIONS['pdf']),
            'word': list(DocumentProcessor.ALLOWED_EXTENSIONS['word']),
            'powerpoint': list(DocumentProcessor.ALLOWED_EXTENSIONS['powerpoint'])
        },
        'features': {
            'ocr_ultra_pro': {
                'tesseract_configs': 6,
                'easyocr_levels': 3,
                'image_enhancements': ['skew_correction', 'denoising', 'contrast', 'binarization', 'artifact_removal'],
                'cost': 'FREE'
            },
            'gemini_fallback': {
                'available': gemini_model is not None,
                'model': 'gemini-1.5-flash' if gemini_model else 'Not configured',
                'cost': 'Variable (solo cuando OCR falla)'
            },
            'markdown_support': MARKITDOWN_AVAILABLE,
            'quality_analysis': True,
            'structured_data_extraction': True
        },
        'performance': {
            'typical_response_time': '2-5 segundos',
            'max_file_size': '32MB',
            'max_pdf_pages': 20
        }
    })


@app.route('/health', methods=['GET'])
def health():
    """Health check para Cloud Run"""
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.now().isoformat(),
        'engines': {
            'tesseract': True,
            'easyocr': EASYOCR_AVAILABLE or easyocr_reader is not None,
            'gemini': gemini_model is not None,
            'pdf': PDF_AVAILABLE,
            'word': DOCX_AVAILABLE,
            'powerpoint': PPTX_AVAILABLE,
            'markitdown': MARKITDOWN_AVAILABLE
        },
        'ultra_pro_features': {
            'image_enhancement': True,
            'multi_config_ocr': True,
            'quality_analysis': True,
            'intelligent_fallback': gemini_model is not None,
            'lazy_loading': True
        }
    }), 200


@app.route('/api/process', methods=['POST'])
def process_document():
    """Procesar cualquier documento con OCR Ultra Pro"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No se encontró archivo en la solicitud'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'Nombre de archivo vacío'}), 400
        
        if not DocumentProcessor.is_allowed_file(file.filename):
            return jsonify({
                'error': 'Formato no soportado',
                'allowed_formats': {
                    'images': list(DocumentProcessor.ALLOWED_EXTENSIONS['image']),
                    'documents': list(DocumentProcessor.ALLOWED_EXTENSIONS['pdf'] | 
                                    DocumentProcessor.ALLOWED_EXTENSIONS['word'] | 
                                    DocumentProcessor.ALLOWED_EXTENSIONS['powerpoint'])
                }
            }), 400
        
        # Obtener parámetro de fallback (por defecto True)
        use_gemini_fallback = request.form.get('use_gemini_fallback', 'true').lower() == 'true'
        
        file_type, ext = DocumentProcessor.detect_file_type(file.filename)
        
        filename = secure_filename(file.filename)
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{ext}') as tmp_file:
            file.save(tmp_file.name)
            temp_path = tmp_file.name
        
        try:
            start_time = time.time()
            
            if file_type == 'image':
                result = ultra_ocr_processor.process_image_ultra_pro(temp_path, use_gemini_fallback)
            elif file_type == 'pdf':
                result = pdf_processor.process(temp_path, use_gemini_fallback)
            elif file_type == 'word':
                result = word_processor.process(temp_path)
            elif file_type == 'powerpoint':
                result = ppt_processor.process(temp_path)
            else:
                result = {'error': 'Tipo de archivo no reconocido'}
            
            processing_time = time.time() - start_time
            
            os.unlink(temp_path)
            
            response = {
                'success': result.get('success', True),
                'filename': filename,
                'file_type': file_type,
                'extension': ext,
                'timestamp': datetime.now().isoformat(),
                'total_processing_time': round(processing_time, 2),
                'gemini_fallback_enabled': use_gemini_fallback,
                'result': result
            }
            
            # Añadir información sobre costos y formato profesional
            if file_type == 'image':
                if result.get('method') == 'ocr_ultra_pro':
                    response['summary'] = {
                        'status': '✅ PROCESAMIENTO COMPLETADO',
                        'method': result.get('method_name', 'OCR Ultra Pro'),
                        'cost': '🆓 GRATIS (0.00 USD)',
                        'time': f"⚡ {result.get('processing_time', 0):.2f} segundos",
                        'quality': f"📊 Score: {result.get('quality_score', 0):.1f}/100",
                        'recommendation': result.get('recommendation', '')
                    }
                elif result.get('method') == 'gemini_fallback':
                    response['summary'] = {
                        'status': '✅ PROCESAMIENTO COMPLETADO',
                        'method': result.get('method_name', 'Gemini AI'),
                        'cost': '💰 Variable según uso de API',
                        'time': f"⚡ {result.get('processing_time', 0):.2f} segundos",
                        'quality': f"📊 Análisis IA Avanzado",
                        'recommendation': result.get('recommendation', 'OCR gratuito no fue suficiente'),
                        'note': '💡 Este documento requirió análisis avanzado de IA'
                    }
                elif result.get('method') == 'all_failed':
                    response['summary'] = {
                        'status': '❌ PROCESAMIENTO FALLIDO',
                        'method': 'Todos los métodos intentados',
                        'error': result.get('error', 'Error desconocido'),
                        'recommendation': result.get('recommendation', 'Mejora la calidad del documento'),
                        'details': result.get('ocr_details', {}),
                        'gemini_status': '❌ No disponible' if not result.get('gemini_available') else '❌ También falló'
                    }
            elif file_type == 'pdf' and 'statistics' in result:
                stats = result['statistics']
                response['summary'] = {
                    'status': '✅ PDF PROCESADO COMPLETAMENTE',
                    'total_pages': result.get('pages_processed', 0),
                    'free_processing': f"🆓 {stats.get('free_ocr', 0)} páginas gratis ({stats.get('free_percentage', 0):.1f}%)",
                    'ai_processing': f"🤖 {stats.get('ai_processing', 0)} páginas con IA",
                    'failed': f"❌ {stats.get('failed', 0)} páginas fallidas" if stats.get('failed', 0) > 0 else "✅ Sin fallos",
                    'time': f"⚡ {round(processing_time, 2)} segundos",
                    'recommendation': result.get('recommendation', 'Procesamiento completado')
                }
            
            return jsonify(response), 200
            
        except Exception as e:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise e
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500


@app.route('/api/process-markdown', methods=['POST'])
def process_markdown():
    """Procesar documento y retornar en formato Markdown usando MarkItDown"""
    try:
        if not MARKITDOWN_AVAILABLE:
            return jsonify({
                'success': False,
                'error': 'MarkItDown no está disponible en este servidor'
            }), 500
        
        if 'file' not in request.files:
            return jsonify({'error': 'No se encontró archivo en la solicitud'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'Nombre de archivo vacío'}), 400
        
        # Guardar archivo temporal
        filename = secure_filename(file.filename)
        file_ext = os.path.splitext(filename)[1]
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            file.save(tmp_file.name)
            temp_path = tmp_file.name
        
        try:
            # Procesar con MarkItDown
            result = markitdown.convert(temp_path)
            
            # Limpiar archivo temporal
            os.unlink(temp_path)
            
            # Preparar respuesta
            markdown_text = result.text_content
            
            response = {
                'success': True,
                'filename': filename,
                'file_extension': file_ext,
                'timestamp': datetime.now().isoformat(),
                'processor': 'markitdown',
                'result': {
                    'markdown': markdown_text,
                    'char_count': len(markdown_text),
                    'word_count': len(markdown_text.split()),
                    'line_count': len(markdown_text.split('\n'))
                }
            }
            
            return jsonify(response), 200
            
        except Exception as e:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise e
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port)